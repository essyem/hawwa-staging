#!/usr/bin/env python3
"""Generate a simple PPTX from PITCH_DECK.md and export PDF for slides and one-pager.

This script creates a basic PPTX with slide titles and bullet points parsed from the markdown.
It also converts `ONE_PAGER.md` to PDF using markdown -> HTML -> wkhtmltopdf if available, otherwise falls back to a simple text-to-PDF.

Note: For production-ready slides, open the generated PPTX in PowerPoint/Google Slides and polish.
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(__file__))
DOCS = os.path.join(ROOT, 'docs', 'startup')
ARTIFACTS = os.path.join(DOCS, 'artifacts')
os.makedirs(ARTIFACTS, exist_ok=True)

PITCH_MD = os.path.join(DOCS, 'PITCH_DECK.md')
ONE_MD = os.path.join(DOCS, 'ONE_PAGER.md')


def parse_markdown_sections(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    # Split into sections by lines that look like 'Slide X' or headings
    sections = []
    current = { 'title': 'Title', 'bullets': [] }
    for line in content.splitlines():
        header = re.match(r'^Slide\s*\d+\s*[-—]?\s*(.*)', line)
        if header:
            if current['title'] or current['bullets']:
                sections.append(current)
            current = {'title': header.group(1).strip() or 'Slide', 'bullets': []}
            continue
        # bullet lines starting with -
        b = re.match(r'^\s*-\s+(.*)', line)
        if b:
            current['bullets'].append(b.group(1).strip())
    if current['title'] or current['bullets']:
        sections.append(current)
    return sections


def make_pptx(sections, out_pptx):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Hawwa Wellness'
    subtitle.text = 'Compassionate postpartum care — Pitch Deck'

    for sec in sections:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = sec['title']
        body = slide.shapes.placeholders[1].text_frame
        for b in sec['bullets']:
            p = body.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)
    prs.save(out_pptx)


def md_to_text(md_path):
    with open(md_path, 'r', encoding='utf-8') as f:
        return f.read()


def text_to_pdf_simple(text, out_pdf):
    # Very small dependency-free PDF writer using reportlab if available
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
    except Exception:
        print('reportlab not available; will write plain text file as fallback')
        with open(out_pdf.replace('.pdf', '.txt'), 'w', encoding='utf-8') as f:
            f.write(text)
        return
    c = canvas.Canvas(out_pdf, pagesize=A4)
    width, height = A4
    margin = 50
    y = height - margin
    for line in text.splitlines():
        if y < margin:
            c.showPage()
            y = height - margin
        c.drawString(margin, y, line[:200])
        y -= 14
    c.save()


if __name__ == '__main__':
    sections = parse_markdown_sections(PITCH_MD)
    pptx_out = os.path.join(ARTIFACTS, 'Hawwa_Pitch_Deck.pptx')
    pdf_out = os.path.join(ARTIFACTS, 'Hawwa_Pitch_Deck.pdf')
    one_pdf = os.path.join(ARTIFACTS, 'Hawwa_One_Pager.pdf')

    print('Generating PPTX:', pptx_out)
    make_pptx(sections, pptx_out)

    # Convert PPTX to PDF: try using libreoffice if available
    if os.system('which libreoffice > /dev/null 2>&1') == 0:
        print('Converting PPTX to PDF with libreoffice')
        os.system(f'libreoffice --headless --convert-to pdf --outdir "{ARTIFACTS}" "{pptx_out}"')
    else:
        print('libreoffice not found; creating a text-PDF fallback for the pitch deck')
        text_to_pdf_simple('\n\n'.join([s['title'] + '\n' + '\n'.join(s['bullets']) for s in sections]), pdf_out)

    # One-pager
    text = md_to_text(ONE_MD)
    text_to_pdf_simple(text, one_pdf)
    print('Artifacts generated in', ARTIFACTS)
